from __future__ import annotations

"""
Create a PowerPoint deck for the R_total zoom-region crossings.

This script reuses the same data-loading and R_total selection logic as the
PRAB zoom-table compiler.  It finds every metric entry inside the configurable
ell--R_total zoom region, resolves the corresponding field-summary PDF in

    D:\PhD\PRAB\figs

and creates a PPTX in

    D:\PhD\PRAB

with one slide per unique PDF.  If a crossing has more than one zoom-region
metric entry, the PDF appears once and the slide table contains multiple rows.

Dependencies
------------
Install these in the Python environment used to run this script:

    pip install python-pptx pymupdf pillow

PyMuPDF is imported as ``fitz`` and is used only to render the first page of
each field-summary PDF to a temporary PNG before insertion into PowerPoint.
"""

import math
import pickle
import re
import tempfile
from collections import defaultdict
from pathlib import Path
from typing import Any, Iterable

import fitz  # PyMuPDF
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


# -----------------------------------------------------------------------------
# Paths and configuration
# -----------------------------------------------------------------------------

F_010_HZ = 1.3e9

DEFAULT_ANALYSIS_ROOT = Path(
    r"D:\PhD\HOMmix\HOMmix_analytical\analysis"
)
DEFAULT_HOMOTYPIC_ROOT = DEFAULT_ANALYSIS_ROOT / "homotypic_rf_multipole"
DEFAULT_HETEROTYPIC_ROOT = DEFAULT_ANALYSIS_ROOT / "heterotypic_crossings"

DEFAULT_PRAB_ROOT = Path(r"D:\PhD\PRAB")
DEFAULT_FIGS_DIR = DEFAULT_PRAB_ROOT / "figs"
DEFAULT_OUTPUT_PPTX = DEFAULT_PRAB_ROOT / "Rtotal_zoom_crossings_field_summaries.pptx"

# These limits should match the ell-versus-R_total zoom figure/table.
ZOOM_ELL_MIN = 0.9
ZOOM_ELL_MAX = 1.1
ZOOM_RTOTAL_MIN = 1.0
ZOOM_RTOTAL_MAX = 6.0

# A row is retained only when max(K_minus, K_plus) is at least the cut-off
# for that row's metric. Set a cut-off to 0.0 to retain every finite result.
DEFAULT_METRIC_CUTOFFS: dict[str, float] = {
    "K_parallel": 0.0,
    "K_perp": 0.0,
    "K_Q": 0.0,
}

# Rendering and slide appearance.
PDF_RENDER_DPI = 220
SLIDE_WIDTH_IN = 13.333
SLIDE_HEIGHT_IN = 7.5

TITLE_FONT_PT = 18
TABLE_FONT_PT = 8
TABLE_HEADER_FONT_PT = 8
FOOTER_FONT_PT = 8

# If True, crossings whose PDFs cannot be found still get slides containing
# their table rows and a "PDF not found" note. If False, they are omitted.
INCLUDE_MISSING_PDFS_AS_SLIDES = True


# -----------------------------------------------------------------------------
# Basic helpers
# -----------------------------------------------------------------------------

def pickle_load(filename: str | Path) -> Any:
    with Path(filename).open("rb") as handle:
        return pickle.load(handle)


def finite_or_nan(value: object) -> float:
    try:
        converted = float(value)
    except Exception:
        return float("nan")
    return converted if math.isfinite(converted) else float("nan")


def abs_finite_or_nan(value: object) -> float:
    try:
        converted = float(abs(complex(value)))
    except Exception:
        return float("nan")
    return converted if math.isfinite(converted) else float("nan")


def first_present(mapping: dict[str, Any], keys: Iterable[str]) -> Any:
    if not isinstance(mapping, dict):
        return float("nan")
    for key in keys:
        if key in mapping:
            return mapping[key]
    return float("nan")


def safe_ratio(numerator: float, denominator: float) -> float:
    numerator = finite_or_nan(numerator)
    denominator = finite_or_nan(denominator)
    if (
        not math.isfinite(numerator)
        or not math.isfinite(denominator)
        or denominator <= 0.0
    ):
        return float("nan")
    return numerator / denominator


def fmt_sci_plain(value: object, significant_figures: int = 3) -> str:
    value = finite_or_nan(value)
    if not math.isfinite(value):
        return "--"
    return f"{value:.{significant_figures - 1}e}"


def fmt_fixed(value: object, decimal_places: int = 4) -> str:
    value = finite_or_nan(value)
    if not math.isfinite(value):
        return "--"
    return f"{value:.{decimal_places}f}"


def normalise_mode_name(mode: object, default_family: str = "TM") -> str:
    text = str(mode).strip()
    if not text or text.lower() == "none":
        return ""

    match = re.search(
        r"(TM|TE)[_\s-]*([0-9]{3,})",
        text,
        flags=re.IGNORECASE,
    )
    if match:
        return f"{match.group(1).upper()}_{match.group(2).zfill(3)}"

    match = re.search(r"\b([0-9]{3,})\b", text)
    if match:
        return f"{default_family.upper()}_{match.group(1).zfill(3)}"

    return text.replace(" ", "_")


def mode_azimuthal_index(mode: object) -> int | None:
    mode = normalise_mode_name(mode)
    match = re.search(r"_(\d)", mode)
    return int(match.group(1)) if match else None


def mode_for_display(mode: object) -> str:
    return normalise_mode_name(mode).replace("_", "")


def metric_for_display(metric_key: str) -> str:
    return {
        "K_parallel": "K_parallel",
        "K_perp": "K_perp",
        "K_Q": "K_Q",
    }[metric_key]


# -----------------------------------------------------------------------------
# Result loading
# -----------------------------------------------------------------------------

def flatten_result_container(data: Any) -> list[dict[str, Any]]:
    results: list[dict[str, Any]] = []

    def visit(value: Any) -> None:
        if isinstance(value, dict):
            if "fields" in value and "crossing" in value:
                results.append(value)
                return
            for nested in value.values():
                visit(nested)
        elif isinstance(value, (list, tuple)):
            for nested in value:
                visit(nested)

    visit(data)
    return results


def load_result_items(
    path_or_root: str | Path,
    *,
    aggregate_filename: str,
    per_folder_filename: str,
) -> list[dict[str, Any]]:
    path = Path(path_or_root)

    if path.is_file():
        return flatten_result_container(pickle_load(path))

    aggregate = path / aggregate_filename
    if aggregate.exists():
        return flatten_result_container(pickle_load(aggregate))

    files = sorted(path.rglob(per_folder_filename))
    if not files:
        raise FileNotFoundError(
            f"No {aggregate_filename!r} or {per_folder_filename!r} found below {path}."
        )

    results: list[dict[str, Any]] = []
    for filename in files:
        results.extend(flatten_result_container(pickle_load(filename)))
    return results


def _result_identity(result: dict[str, Any]) -> tuple[Any, ...]:
    crossing = result.get("crossing", {})
    mode_i, mode_j = result_modes(result)
    family_m = result.get("family_m")
    if family_m is None:
        m_i = mode_azimuthal_index(mode_i)
        m_j = mode_azimuthal_index(mode_j)
        family_m = m_i if m_i == m_j else None

    return (
        family_m,
        mode_i,
        mode_j,
        finite_or_nan(crossing.get("length_factor", float("nan"))),
        finite_or_nan(crossing.get("frequency_Hz", float("nan"))),
    )


def load_homotypic_results(root_or_pickle: str | Path) -> list[dict[str, Any]]:
    path = Path(root_or_pickle)

    if path.is_file():
        return flatten_result_container(pickle_load(path))

    collected: list[dict[str, Any]] = []

    aggregate = path / "all_homotypic_rf_multipole_analyses.pkl"
    if aggregate.exists():
        collected.extend(flatten_result_container(pickle_load(aggregate)))

    per_crossing_files = sorted(path.rglob("homotypic_rf_multipole_analysis.pkl"))
    for filename in per_crossing_files:
        collected.extend(flatten_result_container(pickle_load(filename)))

    if not collected:
        raise FileNotFoundError(
            "No homotypic aggregate or per-crossing RF-multipole results "
            f"were found below {path}."
        )

    unique: dict[tuple[Any, ...], dict[str, Any]] = {}
    for result in collected:
        unique[_result_identity(result)] = result

    results = list(unique.values())
    print(
        "Homotypic loader sources: "
        f"aggregate={'yes' if aggregate.exists() else 'no'}, "
        f"per-crossing files={len(per_crossing_files)}, "
        f"unique crossings={len(results)}"
    )
    return results


def load_heterotypic_results(root_or_pickle: str | Path) -> list[dict[str, Any]]:
    return load_result_items(
        root_or_pickle,
        aggregate_filename="all_heterotypic_multipole_analyses.pkl",
        per_folder_filename="heterotypic_multipole_analysis.pkl",
    )


# -----------------------------------------------------------------------------
# Metric extraction
# -----------------------------------------------------------------------------

METRIC_INFO: dict[str, dict[str, Any]] = {
    "K_parallel": {
        "latex": r"$K_{\parallel}$",
        "explicit_key": "K_parallel_V_per_pC_per_m",
        "legacy_keys": ("loss_like_V_per_pC_per_m",),
    },
    "K_perp": {
        "latex": r"$K_{\perp}$",
        "explicit_key": "K_perp_V_per_pC_per_m2",
        "legacy_keys": (
            "kick_magnitude_V_per_pC_per_m2",
            "kick_mag_V_per_pC_per_m2",
        ),
    },
    "K_Q": {
        "latex": r"$K_Q$",
        "explicit_key": "K_Q_V_per_pC_per_m3",
        "legacy_keys": ("KQ_V_per_pC_per_m3",),
    },
}


def result_modes(result: dict[str, Any]) -> tuple[str, str]:
    crossing = result.get("crossing", {})
    mode_i = normalise_mode_name(result.get("mode_i", crossing.get("mode_i", "")))
    mode_j = normalise_mode_name(result.get("mode_j", crossing.get("mode_j", "")))
    return mode_i, mode_j


def crossing_parameters(result: dict[str, Any]) -> tuple[float, float]:
    crossing = result.get("crossing", {})
    length_factor = finite_or_nan(crossing.get("length_factor", float("nan")))
    frequency_Hz = finite_or_nan(crossing.get("frequency_Hz", float("nan")))
    frequency_normalised = (
        frequency_Hz / F_010_HZ if math.isfinite(frequency_Hz) else float("nan")
    )
    return length_factor, frequency_normalised


def pair_type_key(result: dict[str, Any]) -> str:
    crossing = result.get("crossing", {})
    pair_type = result.get("pair_type") or crossing.get("pair_type")

    if pair_type is None and result.get("crossing_folder"):
        pair_type = Path(result["crossing_folder"]).parent.name

    if pair_type is not None:
        return str(pair_type).lower().replace("-", "_")

    mode_i, mode_j = result_modes(result)
    m_pair = sorted([mode_azimuthal_index(mode_i), mode_azimuthal_index(mode_j)])
    return {
        (0, 1): "monopole_dipole",
        (0, 2): "monopole_quadrupole",
        (1, 2): "dipole_quadrupole",
    }.get(tuple(m_pair), "heterotypic")


def homotypic_family_m(result: dict[str, Any]) -> int:
    family_m = result.get("family_m")
    if family_m is not None:
        return int(family_m)

    mode_i, mode_j = result_modes(result)
    m_i = mode_azimuthal_index(mode_i)
    m_j = mode_azimuthal_index(mode_j)

    if m_i is None or m_j is None or m_i != m_j:
        raise ValueError(f"Could not identify homotypic family for {mode_i}--{mode_j}.")
    return int(m_i)


def figures_of_merit(field_result: dict[str, Any]) -> dict[str, Any]:
    if not isinstance(field_result, dict):
        return {}
    return field_result.get("figures_of_merit", {})


def metric_value(field_result: dict[str, Any], metric_key: str) -> float:
    info = METRIC_INFO[metric_key]
    figures = figures_of_merit(field_result)

    value = abs_finite_or_nan(
        first_present(figures, (info["explicit_key"], *info["legacy_keys"]))
    )
    if math.isfinite(value):
        return value

    if metric_key == "K_parallel":
        value = abs_finite_or_nan(
            field_result.get("kparallel_diagnostics", {})
            .get("fit_V0_U_CST", {})
            .get("k_V_per_pC_per_m", float("nan"))
        )
        if math.isfinite(value):
            return value

    if metric_key == "K_Q":
        Kxx = finite_or_nan(
            first_present(figures, ("K_xx_V_per_pC_per_m3", "Kxx_V_per_pC_per_m3"))
        )
        Kxy = finite_or_nan(
            first_present(figures, ("K_xy_V_per_pC_per_m3", "Kxy_V_per_pC_per_m3"))
        )
        Kyy = finite_or_nan(
            first_present(figures, ("K_yy_V_per_pC_per_m3", "Kyy_V_per_pC_per_m3"))
        )
        if all(math.isfinite(value) for value in (Kxx, Kxy, Kyy)):
            return math.sqrt((Kxx - Kyy) ** 2 + 4.0 * Kxy ** 2)

    return float("nan")


def field_metric_values(result: dict[str, Any], metric_key: str) -> dict[str, float]:
    fields = result.get("fields", {})
    values = {
        field_name: metric_value(fields.get(field_name, {}), metric_key)
        for field_name in ("E1", "E2", "minus", "plus")
    }

    parent_maximum = max(values["E1"], values["E2"])
    mixed_maximum = max(values["minus"], values["plus"])
    parent_total = values["E1"] + values["E2"]
    mixed_total = values["minus"] + values["plus"]

    values["parent_maximum"] = parent_maximum
    values["mixed_maximum"] = mixed_maximum
    values["parent_total"] = parent_total
    values["mixed_total"] = mixed_total
    values["R_total"] = safe_ratio(mixed_total, parent_total)
    return values


def homotypic_metric_key(result: dict[str, Any]) -> str:
    return {
        0: "K_parallel",
        1: "K_perp",
        2: "K_Q",
    }[homotypic_family_m(result)]


def heterotypic_metric_keys(pair_type: str) -> tuple[str, str]:
    mapping = {
        "monopole_dipole": ("K_parallel", "K_perp"),
        "monopole_quadrupole": ("K_parallel", "K_Q"),
        "dipole_quadrupole": ("K_perp", "K_Q"),
    }
    if pair_type not in mapping:
        raise KeyError(f"Unsupported heterotypic pair type: {pair_type}")
    return mapping[pair_type]


def normalised_metric_cutoffs(
    metric_cutoffs: dict[str, float] | None,
) -> dict[str, float]:
    cutoffs = dict(DEFAULT_METRIC_CUTOFFS)
    if metric_cutoffs is not None:
        unknown = set(metric_cutoffs) - set(METRIC_INFO)
        if unknown:
            raise KeyError("Unknown metric cut-off key(s): " + ", ".join(sorted(unknown)))
        cutoffs.update(metric_cutoffs)

    for metric_key, cutoff in cutoffs.items():
        cutoff = finite_or_nan(cutoff)
        if not math.isfinite(cutoff) or cutoff < 0.0:
            raise ValueError(f"Cut-off for {metric_key} must be finite and non-negative.")
        cutoffs[metric_key] = cutoff
    return cutoffs


# -----------------------------------------------------------------------------
# Zoom-region records
# -----------------------------------------------------------------------------

HOMOTYPIC_TYPE_LABELS = {
    0: "M-M",
    1: "D-D",
    2: "Q-Q",
}

HETEROTYPIC_TYPE_LABELS = {
    "monopole_dipole": "M-D",
    "monopole_quadrupole": "M-Q",
    "dipole_quadrupole": "D-Q",
}


def metric_keys_for_record(result: dict[str, Any], *, population: str) -> tuple[str, ...]:
    if population == "homotypic":
        return (homotypic_metric_key(result),)
    if population == "heterotypic":
        return heterotypic_metric_keys(pair_type_key(result))
    raise ValueError(f"Unsupported population: {population!r}")


def zoom_record(
    result: dict[str, Any],
    *,
    population: str,
    metric_key: str,
    cutoffs: dict[str, float],
    ell_min: float,
    ell_max: float,
    rtotal_min: float,
    rtotal_max: float,
) -> dict[str, Any] | None:
    values = field_metric_values(result, metric_key)
    ell, f_hat = crossing_parameters(result)
    r_total = finite_or_nan(values["R_total"])
    mixed_maximum = finite_or_nan(values["mixed_maximum"])

    if (
        not math.isfinite(ell)
        or not math.isfinite(f_hat)
        or not math.isfinite(r_total)
        or not math.isfinite(mixed_maximum)
        or mixed_maximum < cutoffs[metric_key]
        or not (ell_min <= ell <= ell_max)
        or not (rtotal_min <= r_total <= rtotal_max)
    ):
        return None

    mode_i, mode_j = result_modes(result)
    if population == "homotypic":
        crossing_type = HOMOTYPIC_TYPE_LABELS[homotypic_family_m(result)]
        pair_type = None
    else:
        pair_type = pair_type_key(result)
        crossing_type = HETEROTYPIC_TYPE_LABELS[pair_type]

    return {
        "population": population,
        "crossing_type": crossing_type,
        "pair_type": pair_type,
        "mode_i": mode_i,
        "mode_j": mode_j,
        "ell": ell,
        "f_hat": f_hat,
        "metric_key": metric_key,
        "parent_total": values["parent_total"],
        "mixed_total": values["mixed_total"],
        "R_total": r_total,
    }


def collect_zoom_records(
    *,
    homotypic_results: list[dict[str, Any]],
    heterotypic_results: list[dict[str, Any]],
    cutoffs: dict[str, float],
    ell_min: float,
    ell_max: float,
    rtotal_min: float,
    rtotal_max: float,
) -> list[dict[str, Any]]:
    records: list[dict[str, Any]] = []

    for population, results in (
        ("homotypic", homotypic_results),
        ("heterotypic", heterotypic_results),
    ):
        for result in results:
            for metric_key in metric_keys_for_record(result, population=population):
                record = zoom_record(
                    result,
                    population=population,
                    metric_key=metric_key,
                    cutoffs=cutoffs,
                    ell_min=ell_min,
                    ell_max=ell_max,
                    rtotal_min=rtotal_min,
                    rtotal_max=rtotal_max,
                )
                if record is not None:
                    records.append(record)

    return sorted(
        records,
        key=lambda row: (
            row["ell"],
            row["R_total"],
            row["population"],
            row["crossing_type"],
            row["metric_key"],
            row["mode_i"],
            row["mode_j"],
        ),
    )


# -----------------------------------------------------------------------------
# PDF resolution and rendering
# -----------------------------------------------------------------------------

def homotypic_class_key_from_record(record: dict[str, Any]) -> str:
    m_value = mode_azimuthal_index(record["mode_i"])
    return {
        0: "monopole",
        1: "dipole",
        2: "quadrupole",
    }[int(m_value)]


def expected_pdf_name(record: dict[str, Any]) -> str:
    mode_i = record["mode_i"]
    mode_j = record["mode_j"]

    if record["population"] == "homotypic":
        class_key = homotypic_class_key_from_record(record)
        return (
            f"homotypic_{class_key}_"
            f"{mode_i.replace('_', '')}_"
            f"{mode_j.replace('_', '')}_"
            f"field_summary.pdf"
        )

    return (
        f"heterotypic_{record['pair_type']}_"
        f"{mode_i}__{mode_j}__"
        f"field_summary.pdf"
    )


def resolve_pdf_path(record: dict[str, Any], figs_dir: Path) -> Path | None:
    expected = figs_dir / expected_pdf_name(record)
    if expected.exists():
        return expected

    mode_i = record["mode_i"]
    mode_j = record["mode_j"]

    candidate_patterns = []
    if record["population"] == "homotypic":
        class_key = homotypic_class_key_from_record(record)
        candidate_patterns.extend([
            f"homotypic_{class_key}_{mode_i.replace('_', '')}_{mode_j.replace('_', '')}*field_summary.pdf",
            f"homotypic_{class_key}_{mode_i}_{mode_j}*field_summary.pdf",
            f"*{mode_i.replace('_', '')}*{mode_j.replace('_', '')}*field_summary.pdf",
        ])
    else:
        candidate_patterns.extend([
            f"heterotypic_{record['pair_type']}_{mode_i}__{mode_j}__*field_summary.pdf",
            f"heterotypic_{record['pair_type']}*{mode_i}*{mode_j}*field_summary.pdf",
            f"*{mode_i}__{mode_j}__*field_summary.pdf",
        ])

    for pattern in candidate_patterns:
        matches = sorted(figs_dir.glob(pattern))
        if matches:
            return matches[0]

    return None


def render_pdf_first_page_to_png(pdf_path: Path, output_png: Path, *, dpi: int) -> tuple[int, int]:
    document = fitz.open(pdf_path)
    try:
        page = document.load_page(0)
        zoom = dpi / 72.0
        matrix = fitz.Matrix(zoom, zoom)
        pixmap = page.get_pixmap(matrix=matrix, alpha=False)
        pixmap.save(output_png)
        return pixmap.width, pixmap.height
    finally:
        document.close()


def add_picture_contain(
    slide,
    image_path: Path,
    *,
    x_in: float,
    y_in: float,
    w_in: float,
    h_in: float,
    px_width: int,
    px_height: int,
) -> None:
    if px_width <= 0 or px_height <= 0:
        slide.shapes.add_picture(
            str(image_path),
            Inches(x_in),
            Inches(y_in),
            width=Inches(w_in),
            height=Inches(h_in),
        )
        return

    image_aspect = px_width / px_height
    box_aspect = w_in / h_in

    if image_aspect >= box_aspect:
        final_w = w_in
        final_h = w_in / image_aspect
    else:
        final_h = h_in
        final_w = h_in * image_aspect

    final_x = x_in + (w_in - final_w) / 2.0
    final_y = y_in + (h_in - final_h) / 2.0
    slide.shapes.add_picture(
        str(image_path),
        Inches(final_x),
        Inches(final_y),
        width=Inches(final_w),
        height=Inches(final_h),
    )


# -----------------------------------------------------------------------------
# PPTX construction
# -----------------------------------------------------------------------------

def add_textbox(
    slide,
    text: str,
    *,
    x_in: float,
    y_in: float,
    w_in: float,
    h_in: float,
    font_size_pt: float,
    bold: bool = False,
    color: tuple[int, int, int] = (0, 0, 0),
) -> None:
    box = slide.shapes.add_textbox(
        Inches(x_in),
        Inches(y_in),
        Inches(w_in),
        Inches(h_in),
    )
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(0.03)
    frame.margin_right = Inches(0.03)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(font_size_pt)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)


def set_cell_text(cell, text: str, *, font_size_pt: float, bold: bool = False) -> None:
    cell.text = ""
    frame = cell.text_frame
    frame.clear()
    frame.margin_left = Inches(0.03)
    frame.margin_right = Inches(0.03)
    frame.margin_top = Inches(0.01)
    frame.margin_bottom = Inches(0.01)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(font_size_pt)
    run.font.bold = bold


def add_record_table(slide, records: list[dict[str, Any]], *, x_in: float, y_in: float, w_in: float, h_in: float) -> None:
    headers = [
        "Type",
        "Mode 1",
        "Mode 2",
        "ell",
        "fhat",
        "Metric",
        "K1+K2",
        "K+ + K-",
        "Rtotal",
    ]
    rows = len(records) + 1
    cols = len(headers)

    shape = slide.shapes.add_table(
        rows,
        cols,
        Inches(x_in),
        Inches(y_in),
        Inches(w_in),
        Inches(h_in),
    )
    table = shape.table

    # Relative widths tuned for 16:9 slides.
    widths = [0.70, 0.82, 0.82, 0.55, 0.60, 0.85, 1.05, 1.05, 0.65]
    total = sum(widths)
    for index, width in enumerate(widths):
        table.columns[index].width = Inches(w_in * width / total)

    for col_index, header in enumerate(headers):
        set_cell_text(
            table.cell(0, col_index),
            header,
            font_size_pt=TABLE_HEADER_FONT_PT,
            bold=True,
        )

    for row_index, record in enumerate(records, start=1):
        values = [
            record["crossing_type"],
            mode_for_display(record["mode_i"]),
            mode_for_display(record["mode_j"]),
            fmt_fixed(record["ell"], 4),
            fmt_fixed(record["f_hat"], 4),
            metric_for_display(record["metric_key"]),
            fmt_sci_plain(record["parent_total"]),
            fmt_sci_plain(record["mixed_total"]),
            fmt_fixed(record["R_total"], 3),
        ]
        for col_index, value in enumerate(values):
            set_cell_text(
                table.cell(row_index, col_index),
                value,
                font_size_pt=TABLE_FONT_PT,
                bold=False,
            )


def group_records_by_pdf(
    records: list[dict[str, Any]],
    *,
    figs_dir: Path,
) -> tuple[dict[Path | None, list[dict[str, Any]]], list[dict[str, Any]]]:
    grouped: dict[Path | None, list[dict[str, Any]]] = defaultdict(list)
    missing: list[dict[str, Any]] = []

    for record in records:
        pdf_path = resolve_pdf_path(record, figs_dir)
        if pdf_path is None:
            missing.append(record)
            if INCLUDE_MISSING_PDFS_AS_SLIDES:
                grouped[None].append(record)
            continue
        grouped[pdf_path].append(record)

    return grouped, missing


def create_zoom_pptx(
    *,
    homotypic_root_or_pickle: str | Path = DEFAULT_HOMOTYPIC_ROOT,
    heterotypic_root_or_pickle: str | Path = DEFAULT_HETEROTYPIC_ROOT,
    figs_dir: str | Path = DEFAULT_FIGS_DIR,
    output_pptx: str | Path = DEFAULT_OUTPUT_PPTX,
    metric_cutoffs: dict[str, float] | None = None,
    ell_min: float = ZOOM_ELL_MIN,
    ell_max: float = ZOOM_ELL_MAX,
    rtotal_min: float = ZOOM_RTOTAL_MIN,
    rtotal_max: float = ZOOM_RTOTAL_MAX,
    pdf_render_dpi: int = PDF_RENDER_DPI,
) -> Path:
    if ell_min > ell_max:
        raise ValueError("ell_min must not exceed ell_max.")
    if rtotal_min > rtotal_max:
        raise ValueError("rtotal_min must not exceed rtotal_max.")
    if rtotal_min < 0.0:
        raise ValueError("rtotal_min must be non-negative.")

    figs_dir = Path(figs_dir)
    output_pptx = Path(output_pptx)

    cutoffs = normalised_metric_cutoffs(metric_cutoffs)
    homotypic_results = load_homotypic_results(homotypic_root_or_pickle)
    heterotypic_results = load_heterotypic_results(heterotypic_root_or_pickle)

    records = collect_zoom_records(
        homotypic_results=homotypic_results,
        heterotypic_results=heterotypic_results,
        cutoffs=cutoffs,
        ell_min=ell_min,
        ell_max=ell_max,
        rtotal_min=rtotal_min,
        rtotal_max=rtotal_max,
    )

    grouped, missing = group_records_by_pdf(records, figs_dir=figs_dir)

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)
    blank_layout = prs.slide_layouts[6]

    render_dir = output_pptx.parent / "pptx_pdf_render_cache"
    render_dir.mkdir(parents=True, exist_ok=True)

    slide_count = 0
    for pdf_index, (pdf_path, pdf_records) in enumerate(
        sorted(
            grouped.items(),
            key=lambda item: (
                "ZZZ" if item[0] is None else str(item[0]).lower(),
                item[1][0]["ell"],
                item[1][0]["mode_i"],
                item[1][0]["mode_j"],
            ),
        ),
        start=1,
    ):
        if pdf_path is None and not INCLUDE_MISSING_PDFS_AS_SLIDES:
            continue

        first_record = pdf_records[0]
        mode_i = mode_for_display(first_record["mode_i"])
        mode_j = mode_for_display(first_record["mode_j"])
        title = (
            f"{first_record['crossing_type']}  {mode_i} -- {mode_j}  "
            f"(ell={first_record['ell']:.4f}, fhat={first_record['f_hat']:.4f})"
        )

        slide = prs.slides.add_slide(blank_layout)
        slide_count += 1

        add_textbox(
            slide,
            title,
            x_in=0.25,
            y_in=0.12,
            w_in=12.8,
            h_in=0.36,
            font_size_pt=TITLE_FONT_PT,
            bold=True,
        )

        table_height = 0.62 + 0.30 * len(pdf_records)
        table_height = min(table_height, 1.25)
        add_record_table(
            slide,
            pdf_records,
            x_in=0.25,
            y_in=0.58,
            w_in=12.85,
            h_in=table_height,
        )

        image_top = 0.58 + table_height + 0.16
        image_height = 7.15 - image_top

        if pdf_path is None:
            add_textbox(
                slide,
                "Field-summary PDF not found in the configured figs directory.",
                x_in=0.7,
                y_in=image_top + 1.5,
                w_in=12.0,
                h_in=0.6,
                font_size_pt=18,
                bold=True,
                color=(180, 0, 0),
            )
        else:
            output_png = render_dir / f"slide_{slide_count:03d}_{pdf_path.stem}.png"
            px_w, px_h = render_pdf_first_page_to_png(
                pdf_path,
                output_png,
                dpi=pdf_render_dpi,
            )
            add_picture_contain(
                slide,
                output_png,
                x_in=0.25,
                y_in=image_top,
                w_in=12.85,
                h_in=image_height,
                px_width=px_w,
                px_height=px_h,
            )

            add_textbox(
                slide,
                f"PDF: {pdf_path.name}",
                x_in=0.25,
                y_in=7.22,
                w_in=12.85,
                h_in=0.18,
                font_size_pt=FOOTER_FONT_PT,
                color=(90, 90, 90),
            )

    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_pptx)

    print(f"Wrote: {output_pptx}")
    print(f"Zoom metric rows found: {len(records)}")
    print(f"Slides written: {slide_count}")
    print(f"Unique PDFs/slides represented: {len([key for key in grouped if key is not None])}")
    print(f"Missing PDF metric rows: {len(missing)}")
    if missing:
        print("Missing PDFs for:")
        for record in missing:
            print(
                "  "
                f"{record['population']} {record['crossing_type']} "
                f"{record['mode_i']} -- {record['mode_j']} "
                f"{record['metric_key']} "
                f"ell={record['ell']:.4f}, R_total={record['R_total']:.3f}; "
                f"expected {expected_pdf_name(record)}"
            )

    return output_pptx


if __name__ == "__main__":
    create_zoom_pptx(
        homotypic_root_or_pickle=DEFAULT_HOMOTYPIC_ROOT,
        heterotypic_root_or_pickle=DEFAULT_HETEROTYPIC_ROOT,
        figs_dir=DEFAULT_FIGS_DIR,
        output_pptx=DEFAULT_OUTPUT_PPTX,
        metric_cutoffs=DEFAULT_METRIC_CUTOFFS,
        ell_min=ZOOM_ELL_MIN,
        ell_max=ZOOM_ELL_MAX,
        rtotal_min=ZOOM_RTOTAL_MIN,
        rtotal_max=ZOOM_RTOTAL_MAX,
        pdf_render_dpi=PDF_RENDER_DPI,
    )
